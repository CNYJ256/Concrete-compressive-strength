from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape

ROOT = Path(__file__).resolve().parents[2]
PPT_PATH = ROOT / "报告_国际会议版_ACDCB.pptx"


EMU_PER_INCH = 914400


def to_inch(v: int) -> float:
    return v / EMU_PER_INCH


def is_background(shape: BaseShape, slide_w: float, slide_h: float) -> bool:
    w = to_inch(shape.width)
    h = to_inch(shape.height)
    return abs(w - slide_w) < 0.05 and abs(h - slide_h) < 0.05


def is_motif(shape: BaseShape) -> bool:
    w = to_inch(shape.width)
    h = to_inch(shape.height)
    l = to_inch(shape.left)
    # left stripe / tiny dot / tiny arrows etc.
    if l < 0.2 and w <= 0.25:
        return True
    if w <= 0.2 and h <= 0.2:
        return True
    return False


def main() -> None:
    prs = Presentation(str(PPT_PATH))
    slide_w = to_inch(prs.slide_width)
    slide_h = to_inch(prs.slide_height)

    margin_issues: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        for j, shape in enumerate(slide.shapes, start=1):
            if is_background(shape, slide_w, slide_h):
                continue
            if is_motif(shape):
                continue

            left = to_inch(shape.left)
            top = to_inch(shape.top)
            width = to_inch(shape.width)
            height = to_inch(shape.height)
            right = slide_w - (left + width)
            bottom = slide_h - (top + height)

            # Skip footer rows near bottom by design
            if top > 6.95 and height < 0.35:
                continue

            if left < 0.5 or right < 0.5:
                margin_issues.append(
                    f"Slide {i}, shape {j}: horizontal margin too tight (left={left:.2f}, right={right:.2f}, w={width:.2f})"
                )
            if top < 0.3 or bottom < 0.3:
                margin_issues.append(
                    f"Slide {i}, shape {j}: vertical margin too tight (top={top:.2f}, bottom={bottom:.2f}, h={height:.2f})"
                )

    print(f"Slide count: {len(prs.slides)}")
    if margin_issues:
        print("Layout issues detected:")
        for issue in margin_issues:
            print(issue)
    else:
        print("No geometry margin issues detected.")


if __name__ == "__main__":
    main()
