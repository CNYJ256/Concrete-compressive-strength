from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
DATA_PATH = ROOT / "results" / "presentation_work" / "presentation_metrics.json"
OUTPUT_PATH = ROOT / "报告_国际会议版_ACDCB.pptx"
FONT_CN = "Microsoft YaHei"


def rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.strip("#")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))


THEME = {
    "navy": rgb("1E2761"),
    "ice": rgb("CADCFC"),
    "white": rgb("FFFFFF"),
    "light_bg": rgb("F5F7FB"),
    "text": rgb("1F2937"),
    "muted": rgb("5B6470"),
    "accent": rgb("2F5D8A"),
    "accent_2": rgb("5E7FA6"),
    "good": rgb("2E8B57"),
    "warn": rgb("C27A2C"),
    "risk": rgb("B83232"),
}


def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_motif(slide, dark: bool = False) -> None:
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = THEME["ice"] if dark else THEME["accent"]
    stripe.line.fill.background()

    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.06), Inches(0.28), Inches(0.11), Inches(0.11))
    dot.fill.solid()
    dot.fill.fore_color.rgb = THEME["white"] if dark else THEME["ice"]
    dot.line.fill.background()


def add_footer(slide, page: int, total: int, dark: bool = False) -> None:
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.28))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"ACDCB 学术汇报 | 第 {page}/{total} 页"
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.name = FONT_CN
    run.font.size = Pt(11)
    run.font.color.rgb = THEME["ice"] if dark else THEME["muted"]


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.1))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.name = FONT_CN
    run.font.bold = True
    run.font.size = Pt(38)
    run.font.color.rgb = THEME["white"] if dark else THEME["navy"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.25), Inches(11.9), Inches(0.6))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        srun = sp.runs[0]
        srun.font.name = FONT_CN
        srun.font.size = Pt(17)
        srun.font.color.rgb = THEME["ice"] if dark else THEME["muted"]


def add_bullets(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 20,
    dark: bool = False,
) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(9)
        run = p.runs[0]
        run.font.name = FONT_CN
        run.font.size = Pt(font_size)
        run.font.color.rgb = THEME["white"] if dark else THEME["text"]


def add_image(slide, path: Path, x: float, y: float, w: float, h: float, caption: str | None = None) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = THEME["light_bg"]
        ph.line.color.rgb = THEME["muted"]
        tf = ph.text_frame
        tf.text = f"图像缺失：\n{path.name}"
        tf.paragraphs[0].runs[0].font.size = Pt(12)
        tf.paragraphs[0].runs[0].font.name = FONT_CN

    if caption:
        cap = slide.shapes.add_textbox(Inches(x), Inches(y + h + 0.02), Inches(w), Inches(0.26))
        ctf = cap.text_frame
        ctf.clear()
        p = ctf.paragraphs[0]
        p.text = caption
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONT_CN
        run.font.size = Pt(10)
        run.font.color.rgb = THEME["muted"]


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    title_tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.35), Inches(0.4))
    ttf = title_tb.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    tr = tp.runs[0]
    tr.font.name = FONT_CN
    tr.font.bold = True
    tr.font.size = Pt(18)
    tr.font.color.rgb = THEME["white"]

    body_tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.52), Inches(w - 0.35), Inches(h - 0.65))
    btf = body_tb.text_frame
    btf.clear()
    bp = btf.paragraphs[0]
    bp.text = body
    br = bp.runs[0]
    br.font.name = FONT_CN
    br.font.size = Pt(14)
    br.font.color.rgb = THEME["white"]


def add_number_badge(slide, number: str, text: str, x: float, y: float) -> None:
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = THEME["accent"]
    circ.line.fill.background()

    ctf = circ.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = number
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.runs[0]
    cr.font.name = FONT_CN
    cr.font.bold = True
    cr.font.size = Pt(18)
    cr.font.color.rgb = THEME["white"]

    tb = slide.shapes.add_textbox(Inches(x + 0.66), Inches(y + 0.06), Inches(4.8), Inches(0.46))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = FONT_CN
    run.font.size = Pt(15)
    run.font.color.rgb = THEME["text"]


def add_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def make_deck() -> None:
    with DATA_PATH.open("r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_slides = 22
    current = 0

    # Slide 1: Title
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["navy"])
    add_motif(s, dark=True)
    add_title(
        s,
        "ACDCB 混凝土抗压强度预测方法",
        "相对 AdaBoost 基线的系统对比：含消融实验与工程解释",
        dark=True,
    )
    add_bullets(
        s,
        [
            "数据集：UCI Concrete（N=1030，8个输入，1个输出）",
            "评估协议：统一10折交叉验证",
            "汇报风格：国际会议学术技术报告",
        ],
        x=0.75,
        y=2.05,
        w=8.5,
        h=2.2,
        font_size=20,
        dark=True,
    )
    accent = s.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(9.4), Inches(2.0), Inches(3.3), Inches(2.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME["accent"]
    accent.fill.transparency = 22
    accent.line.fill.background()
    add_footer(s, current, total_slides, dark=True)

    # Slide 2: Agenda
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "汇报目录", "从科学问题到工程落地的四段式结构")
    add_number_badge(s, "1", "研究背景与科学问题", 0.9, 2.0)
    add_number_badge(s, "2", "研究方法与数值模拟", 0.9, 2.8)
    add_number_badge(s, "3", "关键实验结果分析", 0.9, 3.6)
    add_number_badge(s, "4", "结论与工程应用", 0.9, 4.4)

    flow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(7.0), Inches(2.2), Inches(5.3), Inches(2.9))
    flow.fill.solid()
    flow.fill.fore_color.rgb = THEME["ice"]
    flow.line.fill.background()
    ftf = flow.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fp.text = "证据驱动\n方法设计\n定量验证\n工程可部署"
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.runs[0]
    fr.font.name = FONT_CN
    fr.font.bold = True
    fr.font.size = Pt(22)
    fr.font.color.rgb = THEME["navy"]
    add_footer(s, current, total_slides)

    # Slide 3: Background
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "背景：为何抗压强度预测仍是核心问题")
    add_bullets(
        s,
        [
            "抗压强度直接决定结构安全、质量控制与全寿命评估。",
            "实体试验可靠但成本高、周期长，不利于配合比快速迭代。",
            "数据驱动模型可显著加速材料筛选与设计决策。",
            "核心挑战在于胶凝体系、用水、骨料与龄期之间的强非线性耦合。",
        ],
        x=0.75,
        y=1.8,
        w=6.1,
        h=4.8,
        font_size=18,
    )
    add_image(
        s,
        ROOT / data["existing_figures"]["data_distribution"],
        x=6.9,
        y=1.8,
        w=5.9,
        h=4.7,
        caption="UCI 数据分布（龄期采用对数坐标）",
    )
    add_footer(s, current, total_slides)

    # Slide 4: Scientific gap
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "现有报告与文献中的关键科学缺口")
    add_card(
        s,
        "非线性表征不足",
        "单模型假设难以充分刻画配比参数与龄期的耦合关系。",
        0.8,
        1.8,
        3.9,
        2.0,
        THEME["accent"],
    )
    add_card(
        s,
        "分布迁移",
        "早龄期与后龄期样本呈现不同统计规律。",
        4.95,
        1.8,
        3.9,
        2.0,
        THEME["accent_2"],
    )
    add_card(
        s,
        "模型偏好漂移",
        "不同学习器在不同龄期与特征子空间中优势并不一致。",
        9.1,
        1.8,
        3.4,
        2.0,
        THEME["navy"],
    )
    add_bullets(
        s,
        [
            "原始 报告.pptx 已识别上述痛点，但结果页存在不完整内容。",
            "本版本以论文证据、可复现实验指标与高分辨率图表补齐关键结论。",
        ],
        x=0.8,
        y=4.2,
        w=12.0,
        h=2.0,
        font_size=18,
    )
    add_footer(s, current, total_slides)

    # Slide 5: Literature synthesis
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "文献脉络：paper2 → paper1 → ACDCB")

    table = s.shapes.add_table(4, 4, Inches(0.8), Inches(1.85), Inches(12.0), Inches(3.4)).table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3.6)
    table.columns[3].width = Inches(3.7)

    headers = ["工作", "模型家族", "优势", "局限 / 下一步"]
    rows = [
        ["Yeh 1998（paper2）", "ANN + 回归", "验证 ANN 优于传统回归", "缺少双空间结构，鲁棒性有限"],
        ["Feng 2020（paper1）", "AdaBoost", "集成学习基线表现强", "单空间建模，无龄期分段融合"],
        ["本研究", "ACDCB", "双空间 + 约束融合 + 龄期分段", "仍需外部多场景验证"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["navy"]
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = THEME["white"]
        p.runs[0].font.size = Pt(14)

    for r in range(1, 4):
        for c in range(4):
            cell = table.cell(r, c)
            cell.text = rows[r - 1][c]
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = FONT_CN
            p.runs[0].font.size = Pt(13)
            p.runs[0].font.color.rgb = THEME["text"]

    add_bullets(
        s,
        ["方法演进主线：从单一预测器走向结构化的双空间约束融合。"],
        x=0.9,
        y=5.5,
        w=11.8,
        h=1.0,
        font_size=16,
    )
    add_footer(s, current, total_slides)

    # Slide 6: Dataset and protocol
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "数据集与评估协议")

    add_card(s, "1030", "样本数", 0.8, 1.9, 2.2, 1.4, THEME["navy"])
    add_card(s, "8", "输入变量", 3.2, 1.9, 2.2, 1.4, THEME["accent"])
    add_card(s, "1", "目标（MPa）", 5.6, 1.9, 2.2, 1.4, THEME["accent_2"])
    add_card(s, "10折", "主评估协议", 8.0, 1.9, 2.4, 1.4, THEME["navy"])

    add_bullets(
        s,
        [
            "评价指标：$R^2$、RMSE、MAE、MAPE（统一协议确保公平对比）。",
            "龄期分布偏态明显（skew = 3.269），支持采用对数可视化与龄期分段。",
            "目标强度范围 2.33–82.60 MPa，要求模型在低/高强区间均具鲁棒性。",
        ],
        x=0.8,
        y=3.65,
        w=7.1,
        h=2.4,
        font_size=16,
    )

    formula_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.05), Inches(3.55), Inches(4.75), Inches(2.2))
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = THEME["ice"]
    formula_box.line.fill.background()
    tf = formula_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "协议目标：\n在统一折分下最大化 $R^2$，同时最小化 RMSE/MAE/MAPE"
    p.runs[0].font.name = FONT_CN
    p.runs[0].font.size = Pt(17)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = THEME["navy"]
    add_footer(s, current, total_slides)

    # Slide 7: Baseline reproducibility
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "基线复现快照（paper1设置）")

    table = s.shapes.add_table(4, 5, Inches(0.8), Inches(1.95), Inches(6.8), Inches(3.9)).table
    headers = ["模型", "$R^2$", "RMSE", "MAE", "MAPE"]
    baseline_rows = [
        ["AdaBoost（10折）", "0.9090", "4.9695", "3.5085", "13.3513"],
        ["ANN（9:1划分）", "0.9160", "4.7273", "3.3637", "11.91"],
        ["SVM（9:1划分）", "0.8713", "5.8492", "4.2099", "14.73"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["navy"]
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = THEME["white"]

    for r in range(1, 4):
        for c in range(5):
            cell = table.cell(r, c)
            cell.text = baseline_rows[r - 1][c]
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.name = FONT_CN
            p.runs[0].font.size = Pt(13)

    add_image(
        s,
        ROOT / "figures" / "presentation_highres" / "04_model_family_scatter.png",
        x=7.85,
        y=1.95,
        w=4.9,
        h=3.9,
        caption="跨模型误差-精度分布图",
    )
    add_footer(s, current, total_slides)

    # Slide 8: ACDCB architecture
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "ACDCB 架构：龄期分段双空间约束融合")

    # Diagram boxes
    boxes = [
        (0.8, 2.1, 2.2, 1.0, "输入层\n8个原始变量"),
        (3.3, 1.4, 2.5, 1.0, "主空间\n（32维特征）"),
        (3.3, 2.8, 2.5, 1.0, "锚空间\n（22维特征）"),
        (6.2, 2.1, 2.7, 1.0, "模型池\nXGB/LGBM/HGB"),
        (9.2, 2.1, 3.0, 1.0, "约束融合\n$w_i\\ge0,\\sum w_i=1$"),
    ]

    for x, y, w, h, text in boxes:
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        b.fill.solid()
        b.fill.fore_color.rgb = THEME["ice"]
        b.line.color.rgb = THEME["accent"]
        tf = b.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT_CN
        r.font.bold = True
        r.font.size = Pt(14)
        r.font.color.rgb = THEME["navy"]

    for x in [2.95, 5.95, 8.95]:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.35), Inches(0.28), Inches(0.46))
        arr.fill.solid()
        arr.fill.fore_color.rgb = THEME["accent"]
        arr.line.fill.background()

    split = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.6), Inches(3.4), Inches(0.4), Inches(0.9))
    split.fill.solid()
    split.fill.fore_color.rgb = THEME["accent_2"]
    split.line.fill.background()

    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.25), Inches(4.45), Inches(3.0), Inches(1.2))
    tag.fill.solid()
    tag.fill.fore_color.rgb = THEME["navy"]
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "龄期分段输出\nage≤28: $w_e$\nage>28: $w_l$"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = FONT_CN
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = THEME["white"]
    add_footer(s, current, total_slides)

    # Slide 9: Dual-space features
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "双空间特征设计")

    add_card(
        s,
        "共享机制特征（14项）",
        "binder、w/c、w/b、age_log1p、age_sqrt、age^0.25、abrams_index，\ncement-age interaction、paste index 等",
        0.8,
        1.9,
        5.9,
        2.3,
        THEME["accent"],
    )
    add_card(
        s,
        "主空间增强特征（10项）",
        "binder/agg、water/paste、cement_fraction、slag_fraction，\nflyash_fraction、maturity_index、age_inverse 等",
        7.0,
        1.9,
        5.5,
        2.3,
        THEME["navy"],
    )

    eq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.45), Inches(11.7), Inches(1.6))
    eq.fill.solid()
    eq.fill.fore_color.rgb = THEME["ice"]
    eq.line.fill.background()
    tf = eq.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "$\\mathbf{z}^{(p)}=\\phi_p(\\mathbf{x}),\\; \\mathbf{z}^{(a)}=\\phi_a(\\mathbf{x})$    |    锚空间：22维，主空间：32维"
    r = p.runs[0]
    r.font.name = FONT_CN
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = THEME["navy"]

    p2 = tf.add_paragraph()
    p2.text = "设计意图：高容量表达（主空间）+ 稳健正则（锚空间）。"
    p2.runs[0].font.name = FONT_CN
    p2.runs[0].font.size = Pt(14)
    p2.runs[0].font.color.rgb = THEME["text"]
    add_footer(s, current, total_slides)

    # Slide 10: constrained optimization
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "约束融合优化公式")

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(7.4), Inches(3.9))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME["ice"]
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "$\\min_{\\mathbf{w}}\\;RMSE(\\mathbf{y},\\mathbf{P}\\mathbf{w})$"
    p.runs[0].font.name = FONT_CN
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = THEME["navy"]

    p2 = tf.add_paragraph()
    p2.text = "约束条件：$w_i\\ge0$ 且 $\\sum_i w_i=1$"
    p2.runs[0].font.name = FONT_CN
    p2.runs[0].font.size = Pt(24)
    p2.runs[0].font.bold = True
    p2.runs[0].font.color.rgb = THEME["accent"]

    p3 = tf.add_paragraph()
    p3.text = "单纯形约束可避免不稳定的负权抵消，并保持权重可解释性。"
    p3.runs[0].font.name = FONT_CN
    p3.runs[0].font.size = Pt(15)
    p3.runs[0].font.color.rgb = THEME["text"]

    add_card(
        s,
        "方法价值",
        "1）优化过程稳定\n2）凸组合解释清晰\n3）跨折对比可复现",
        8.45,
        1.8,
        4.1,
        3.9,
        THEME["navy"],
    )
    add_footer(s, current, total_slides)

    # Slide 11: age-conditioned
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "龄期分段融合策略")
    add_bullets(
        s,
        [
            "阈值设为28天：符合工程习惯，同时对应明显的分布切换点。",
            "$\\mathbf{w}_e$ 在早龄期子集优化，$\\mathbf{w}_l$ 在晚龄期子集优化。",
            "代码中的选择规则：优先比较 $R^2$，并以 RMSE 作为平局判据（$\\tau=5\\times10^{-4}$）。",
        ],
        x=0.8,
        y=1.9,
        w=6.2,
        h=2.5,
        font_size=16,
    )

    add_image(
        s,
        ROOT / "figures" / "presentation_highres" / "03_age_segment_weights.png",
        x=7.0,
        y=1.85,
        w=5.7,
        h=4.7,
        caption="不同龄期下的权重迁移",
    )
    add_footer(s, current, total_slides)

    # Slide 12: numerical simulation settings
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "数值实验设置")

    table = s.shapes.add_table(5, 2, Inches(0.85), Inches(1.85), Inches(7.1), Inches(4.6)).table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(4.8)

    hp_rows = [
        ["XGBoost", "n_estimators=1482, lr=0.0446, depth=4, gamma=2.4344"],
        ["LightGBM", "n_estimators=2127, lr=0.0320, leaves=32, depth=4"],
        ["HGB（主空间）", "max_iter=1809, lr=0.0568, depth=12, max_bins=213"],
        ["HGB（锚空间）", "max_iter=2400, lr=0.0280, max_leaf_nodes=15"],
        ["优化器", "在单纯形约束下采用 SLSQP 求解全局与分段权重"],
    ]

    for r in range(5):
        for c in range(2):
            cell = table.cell(r, c)
            cell.text = hp_rows[r][c]
            if c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME["navy"]
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = THEME["white"]
                run.font.bold = True
                run.font.name = FONT_CN
            else:
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.color.rgb = THEME["text"]
                run.font.name = FONT_CN
            run.font.size = Pt(13)

    opt = data["optimizer_summary"]
    add_card(
        s,
        "收敛证据",
        (
            f"早龄期 RMSE：{opt['early']['start_rmse']:.4f} → {opt['early']['end_rmse']:.4f}"
            f"（{opt['early']['iterations']}次迭代）\n"
            f"晚龄期 RMSE：{opt['late']['start_rmse']:.4f} → {opt['late']['end_rmse']:.4f}"
            f"（{opt['late']['iterations']}次迭代）"
        ),
        8.2,
        2.2,
        4.2,
        2.6,
        THEME["accent"],
    )
    add_footer(s, current, total_slides)

    # Slide 13: main protocol results
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "关键结果 I：统一10折协议对比")
    add_image(
        s,
        ROOT / "figures" / "presentation_highres" / "01_main_metrics_2x2.png",
        x=0.8,
        y=1.55,
        w=12.0,
        h=5.55,
        caption="相较 paper1 AdaBoost，ACDCB 在四项指标上均稳定提升",
    )
    add_footer(s, current, total_slides)

    # Slide 14: true vs pred
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "关键结果 II：OOF 真值-预测分布")
    add_image(
        s,
        ROOT / data["existing_figures"]["true_vs_pred"],
        x=0.75,
        y=1.7,
        w=8.6,
        h=5.3,
        caption="六边形密度图（对数尺度）：ACDCB 点云更集中于 y=x 附近",
    )
    add_bullets(
        s,
        [
            "采用 hexbin + 对数密度可视化，有效避免样本重叠遮挡。",
            "误差带收窄与 RMSE 从 4.9695 降至 3.6995 MPa 的结果一致。",
            "点云与对角线贴合度提升，说明全强度区间泛化能力更强。",
        ],
        x=9.55,
        y=2.0,
        w=3.25,
        h=4.4,
        font_size=14,
    )
    add_footer(s, current, total_slides)

    # Slide 15: ablation zoomed
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "关键结果 III：近值可读的消融分析")
    add_image(
        s,
        ROOT / "figures" / "presentation_highres" / "02_ablation_zoomed.png",
        x=0.75,
        y=1.7,
        w=11.95,
        h=4.9,
        caption="通过坐标轴放大，让亚百分点级提升可见且可审计",
    )
    add_bullets(
        s,
        [
            "V2 相比 V1：双空间锚定带来稳定正向增益。",
            "V3 相比 V2：龄期分段进一步带来小幅且可复现提升。",
            "V4（raw）揭示指标间权衡关系，提示后续联合寻优方向。",
        ],
        x=0.9,
        y=6.25,
        w=11.8,
        h=0.8,
        font_size=13,
    )
    add_footer(s, current, total_slides)

    # Slide 16: stability visuals
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "跨折稳定性证据")
    add_image(
        s,
        ROOT / data["existing_figures"]["ablation_r2_rmse"],
        x=0.8,
        y=1.9,
        w=5.95,
        h=4.8,
        caption="消融实验的 R²/RMSE 分布图",
    )
    add_image(
        s,
        ROOT / data["existing_figures"]["fold_r2_boxplot"],
        x=6.85,
        y=1.9,
        w=5.95,
        h=4.8,
        caption="折级别 R² 分布抬升",
    )
    add_footer(s, current, total_slides)

    # Slide 17: optimizer convergence
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "约束优化器收敛过程")
    add_image(
        s,
        ROOT / data["existing_figures"]["optimizer_convergence"],
        x=0.8,
        y=1.85,
        w=6.4,
        h=4.9,
        caption="SLSQP 在早/晚龄期分段下的目标函数收敛轨迹",
    )

    add_card(
        s,
        "结果解读",
        "两段优化均在10–13次迭代内单调收敛。\n未观察到发散或明显震荡。",
        7.5,
        2.0,
        5.0,
        1.9,
        THEME["accent"],
    )
    add_card(
        s,
        "工程意义",
        "求解器稳定性满足常规重训练需求，支撑可复现的模型治理流程。",
        7.5,
        4.1,
        5.0,
        1.9,
        THEME["navy"],
    )
    add_footer(s, current, total_slides)

    # Slide 18: ann/svm comparison
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "与 ANN/SVM 及基线集成模型对比")
    add_image(
        s,
        ROOT / "figures" / "presentation_highres" / "04_model_family_scatter.png",
        x=0.8,
        y=1.8,
        w=7.7,
        h=5.1,
        caption="ACDCB 位于更优的 R²-RMSE 组合区间",
    )

    comp_table = s.shapes.add_table(4, 3, Inches(8.7), Inches(1.95), Inches(3.95), Inches(3.6)).table
    comp_table.columns[0].width = Inches(1.7)
    comp_table.columns[1].width = Inches(1.0)
    comp_table.columns[2].width = Inches(1.2)

    rows = [
        ["模型", "R²", "RMSE"],
        ["ACDCB", "0.9488", "3.6995"],
        ["ANN", "0.9160", "4.7273"],
        ["SVM", "0.8713", "5.8492"],
    ]

    for r in range(4):
        for c in range(3):
            cell = comp_table.cell(r, c)
            cell.text = rows[r][c]
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = THEME["navy"]
                run.font.color.rgb = THEME["white"]
                run.font.bold = True
                run.font.name = FONT_CN
            else:
                run.font.color.rgb = THEME["text"]
                run.font.name = FONT_CN
                if r == 1:
                    run.font.bold = True
            run.font.size = Pt(13)

    add_bullets(
        s,
        ["在更严格的统一10折协议下，结果仍明显优于单次随机划分基线。"],
        x=8.75,
        y=5.8,
        w=4.0,
        h=0.9,
        font_size=12,
    )
    add_footer(s, current, total_slides)

    # Slide 19: engineering interpretation
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "性能提升的工程解释")

    add_card(
        s,
        "精度收益",
        "$\\Delta R^2=+0.0398$，RMSE 下降 1.2699 MPa，可提升安全裕度评估可靠性。",
        0.9,
        1.9,
        3.9,
        2.2,
        THEME["accent"],
    )
    add_card(
        s,
        "稳定性收益",
        "约束融合降低跨折波动，并避免负权重引发的不稳定抵消。",
        4.95,
        1.9,
        3.9,
        2.2,
        THEME["accent_2"],
    )
    add_card(
        s,
        "可解释性收益",
        "分段权重可直观揭示不同龄期下模型主导关系（早龄期 vs 晚龄期）。",
        9.0,
        1.9,
        3.4,
        2.2,
        THEME["navy"],
    )

    add_bullets(
        s,
        [
            "这不仅是预测精度提升，更是面向材料工程师的决策支持能力提升。",
            "权重轨迹具备可审计性，可与领域专家进行机理层面的联合讨论。",
        ],
        x=0.9,
        y=4.5,
        w=11.9,
        h=1.7,
        font_size=18,
    )
    add_footer(s, current, total_slides)

    # Slide 20: engineering applications
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "工程应用路径")

    stages = [
        ("配合比筛选", "在实验前快速完成候选方案排序"),
        ("试验规划", "减少重复破坏性试验，降低成本"),
        ("现场质控", "对欠强风险批次进行实时预警"),
        ("全寿命管理", "基于龄期进行强度预测与维护决策"),
    ]

    x0 = 0.9
    for i, (name, desc) in enumerate(stages):
        x = x0 + i * 3.05
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(2.6), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = THEME["ice"] if i % 2 == 0 else THEME["light_bg"]
        card.line.color.rgb = THEME["accent"]

        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = name
        p.runs[0].font.name = FONT_CN
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(17)
        p.runs[0].font.color.rgb = THEME["navy"]

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.runs[0].font.name = FONT_CN
        p2.runs[0].font.size = Pt(13)
        p2.runs[0].font.color.rgb = THEME["text"]

        if i < len(stages) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.65), Inches(3.3), Inches(0.35), Inches(0.45))
            arr.fill.solid()
            arr.fill.fore_color.rgb = THEME["accent"]
            arr.line.fill.background()

    add_footer(s, current, total_slides)

    # Slide 21: limitations
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["light_bg"])
    add_motif(s)
    add_title(s, "局限性、风险与应对")

    matrix = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.9), Inches(7.4), Inches(4.6))
    matrix.fill.solid()
    matrix.fill.fore_color.rgb = THEME["ice"]
    matrix.line.fill.background()

    add_bullets(
        s,
        [
            "V4 raw 方案在 R²/RMSE 上优于 V3，但 MAE 未同步改善，仍存在多指标权衡。",
            "当前结论主要来自单一公开数据集，仍需外部多区域数据验证。",
            "后续工作：工程特征与原始特征流程独立调优，并补充 SHAP/PDP 解释分析。",
        ],
        x=1.15,
        y=2.2,
        w=6.9,
        h=4.0,
        font_size=16,
    )

    add_card(
        s,
        "风险等级",
        "数据迁移风险：中\n模型不稳定风险：低\n可解释性缺口：中",
        8.7,
        2.0,
        3.9,
        1.9,
        THEME["warn"],
    )
    add_card(
        s,
        "应对策略",
        "跨数据集验证\n联合超参数搜索\n周期性重校准与监控",
        8.7,
        4.1,
        3.9,
        1.9,
        THEME["good"],
    )
    add_footer(s, current, total_slides)

    # Slide 22: conclusion
    current += 1
    s = add_slide(prs)
    set_slide_bg(s, THEME["navy"])
    add_motif(s, dark=True)
    add_title(s, "结论与工程展望", dark=True)

    add_bullets(
        s,
        [
            "在统一10折协议下，ACDCB 相对 paper1 AdaBoost 实现了稳定提升。",
            "性能提升来源可追踪：双空间锚定 + 龄期分段约束融合。",
            "方法已具备工程部署潜力，可服务于配合比设计与质量控制决策。",
            "下一步：外部验证、阈值迁移性研究与可解释性增强。",
        ],
        x=0.8,
        y=1.9,
        w=8.8,
        h=3.6,
        font_size=19,
        dark=True,
    )

    closing = s.shapes.add_shape(MSO_SHAPE.SUN, Inches(9.8), Inches(2.2), Inches(2.8), Inches(2.8))
    closing.fill.solid()
    closing.fill.fore_color.rgb = THEME["ice"]
    closing.fill.transparency = 18
    closing.line.fill.background()

    ctf = closing.text_frame
    ctf.clear()
    cp = ctf.paragraphs[0]
    cp.text = "谢谢！"
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.runs[0]
    cr.font.name = FONT_CN
    cr.font.bold = True
    cr.font.size = Pt(30)
    cr.font.color.rgb = THEME["navy"]

    add_footer(s, current, total_slides, dark=True)

    prs.save(str(OUTPUT_PATH))
    print(f"PPT已生成：{OUTPUT_PATH}")


if __name__ == "__main__":
    make_deck()
